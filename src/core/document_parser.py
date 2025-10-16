"""
文档解析模块
支持本地解析PPT、Excel、Word等文件，提取文本内容
文件不会上传到网络，仅提取文本信息用于AI分析
"""

import os
import json
import csv
from typing import Dict, List, Any
from pathlib import Path
from datetime import datetime


class DocumentParser:
    """文档解析器基类"""

    @staticmethod
    def extract_text_from_excel(file_path: str) -> Dict[str, Any]:
        """
        从Excel文件提取文本内容

        Args:
            file_path: Excel文件路径

        Returns:
            提取的内容字典
        """
        try:
            import openpyxl
        except ImportError:
            raise ImportError("请安装 openpyxl: pip install openpyxl")

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        workbook = openpyxl.load_workbook(file_path, data_only=True)
        result = {
            "file_name": os.path.basename(file_path),
            "file_type": "excel",
            "sheets": {}
        }

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_data = []

            for row in sheet.iter_rows(values_only=True):
                # 过滤空行
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(row_data):  # 如果行中有非空内容
                    sheet_data.append(row_data)

            result["sheets"][sheet_name] = sheet_data

        print(f"✅ 已提取Excel内容: {file_path}")
        return result

    @staticmethod
    def extract_text_from_word(file_path: str) -> Dict[str, Any]:
        """
        从Word文件提取文本内容

        Args:
            file_path: Word文件路径

        Returns:
            提取的内容字典
        """
        try:
            from docx import Document
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        doc = Document(file_path)
        result = {
            "file_name": os.path.basename(file_path),
            "file_type": "word",
            "paragraphs": [],
            "tables": []
        }

        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                result["paragraphs"].append(para.text)

        # 提取表格
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            result["tables"].append(table_data)

        print(f"✅ 已提取Word内容: {file_path}")
        return result

    @staticmethod
    def extract_text_from_ppt(file_path: str) -> Dict[str, Any]:
        """
        从PowerPoint文件提取文本内容

        Args:
            file_path: PPT文件路径

        Returns:
            提取的内容字典
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        prs = Presentation(file_path)
        result = {
            "file_name": os.path.basename(file_path),
            "file_type": "powerpoint",
            "slides": []
        }

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = {
                "slide_number": slide_num,
                "texts": [],
                "notes": ""
            }

            # 提取幻灯片中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content["texts"].append(shape.text.strip())

            # 提取备注
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    slide_content["notes"] = notes_frame.text.strip()

            result["slides"].append(slide_content)

        print(f"✅ 已提取PPT内容: {file_path}")
        return result

    @staticmethod
    def extract_text_from_pdf(file_path: str) -> Dict[str, Any]:
        """
        从PDF文件提取文本内容

        Args:
            file_path: PDF文件路径

        Returns:
            提取的内容字典
        """
        try:
            import PyPDF2
        except ImportError:
            raise ImportError("请安装 PyPDF2: pip install PyPDF2")

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        result = {
            "file_name": os.path.basename(file_path),
            "file_type": "pdf",
            "pages": []
        }

        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)

            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    result["pages"].append({
                        "page_number": page_num,
                        "text": text.strip()
                    })

        print(f"✅ 已提取PDF内容: {file_path}")
        return result

    @classmethod
    def parse_document(cls, file_path: str) -> Dict[str, Any]:
        """
        自动识别文件类型并解析

        Args:
            file_path: 文件路径

        Returns:
            提取的内容字典
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        file_ext = os.path.splitext(file_path)[1].lower()

        parsers = {
            '.xlsx': cls.extract_text_from_excel,
            '.xls': cls.extract_text_from_excel,
            '.docx': cls.extract_text_from_word,
            '.doc': cls.extract_text_from_word,
            '.pptx': cls.extract_text_from_ppt,
            '.ppt': cls.extract_text_from_ppt,
            '.pdf': cls.extract_text_from_pdf,
        }

        parser = parsers.get(file_ext)
        if parser is None:
            raise ValueError(f"不支持的文件格式: {file_ext}")

        return parser(file_path)

    @staticmethod
    def format_extracted_data(data: Dict[str, Any]) -> str:
        """
        将提取的数据格式化为文本字符串，用于AI分析

        Args:
            data: 提取的文档数据

        Returns:
            格式化的文本字符串
        """
        file_type = data.get("file_type", "unknown")
        file_name = data.get("file_name", "unknown")

        output = [f"文件名: {file_name}", f"文件类型: {file_type}", "=" * 50, ""]

        if file_type == "excel":
            for sheet_name, sheet_data in data.get("sheets", {}).items():
                output.append(f"\n工作表: {sheet_name}")
                output.append("-" * 40)
                for row in sheet_data:
                    output.append(" | ".join(row))

        elif file_type == "word":
            output.append("段落内容:")
            output.append("-" * 40)
            for para in data.get("paragraphs", []):
                output.append(para)
                output.append("")

            if data.get("tables"):
                output.append("\n表格内容:")
                output.append("-" * 40)
                for i, table in enumerate(data.get("tables", []), 1):
                    output.append(f"\n表格 {i}:")
                    for row in table:
                        output.append(" | ".join(row))

        elif file_type == "powerpoint":
            for slide in data.get("slides", []):
                slide_num = slide.get("slide_number", "?")
                output.append(f"\n幻灯片 {slide_num}:")
                output.append("-" * 40)
                for text in slide.get("texts", []):
                    output.append(text)
                if slide.get("notes"):
                    output.append(f"\n备注: {slide['notes']}")
                output.append("")

        elif file_type == "pdf":
            for page in data.get("pages", []):
                page_num = page.get("page_number", "?")
                output.append(f"\n第 {page_num} 页:")
                output.append("-" * 40)
                output.append(page.get("text", ""))
                output.append("")

        return "\n".join(output)


def parse_multiple_documents(file_paths: List[str]) -> Dict[str, str]:
    """
    批量解析多个文档

    Args:
        file_paths: 文件路径列表

    Returns:
        文件路径到格式化文本的映射
    """
    results = {}

    for file_path in file_paths:
        try:
            data = DocumentParser.parse_document(file_path)
            formatted_text = DocumentParser.format_extracted_data(data)
            results[file_path] = formatted_text
        except Exception as e:
            print(f"❌ 解析失败 {file_path}: {str(e)}")
            results[file_path] = f"解析失败: {str(e)}"

    return results


class DataExporter:
    """数据导出器 - 支持JSON、CSV等格式"""

    @staticmethod
    def save_as_json(data: Dict[str, Any], output_path: str, pretty: bool = True) -> str:
        """
        将提取的数据保存为JSON文件

        Args:
            data: 提取的文档数据字典
            output_path: 输出文件路径
            pretty: 是否格式化输出(方便阅读)

        Returns:
            保存的文件路径
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # 添加元数据
        export_data = {
            "metadata": {
                "exported_at": datetime.now().isoformat(),
                "exporter": "SellSysInsurance DataExporter",
                "version": "1.0"
            },
            "data": data
        }

        with open(output_path, 'w', encoding='utf-8') as f:
            if pretty:
                json.dump(export_data, f, ensure_ascii=False, indent=2)
            else:
                json.dump(export_data, f, ensure_ascii=False)

        print(f"✅ JSON文件已保存: {output_path}")
        return str(output_path)

    @staticmethod
    def save_as_csv(data: Dict[str, Any], output_path: str) -> str:
        """
        将Excel类型的数据保存为CSV文件

        Args:
            data: 提取的文档数据字典
            output_path: 输出文件路径

        Returns:
            保存的文件路径(如果是Excel数据会保存多个CSV文件)

        Note:
            仅支持Excel类型的数据,其他类型建议使用JSON
        """
        file_type = data.get("file_type", "unknown")

        if file_type != "excel":
            raise ValueError(f"CSV导出仅支持Excel文件,当前文件类型: {file_type}")

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # 如果有多个sheet,为每个sheet创建单独的CSV文件
        sheets = data.get("sheets", {})
        saved_files = []

        if len(sheets) == 1:
            # 只有一个sheet,直接保存为指定文件名
            sheet_name, sheet_data = list(sheets.items())[0]
            DataExporter._write_csv(sheet_data, output_path)
            saved_files.append(str(output_path))
        else:
            # 多个sheet,为每个创建单独文件
            base_name = output_path.stem
            base_dir = output_path.parent

            for sheet_name, sheet_data in sheets.items():
                # 清理sheet名称,去除特殊字符
                safe_sheet_name = "".join(c if c.isalnum() or c in "._- " else "_" for c in sheet_name)
                csv_path = base_dir / f"{base_name}_{safe_sheet_name}.csv"
                DataExporter._write_csv(sheet_data, csv_path)
                saved_files.append(str(csv_path))

        print(f"✅ CSV文件已保存: {len(saved_files)}个文件")
        for f in saved_files:
            print(f"   - {f}")

        return saved_files[0] if len(saved_files) == 1 else str(output_path.parent)

    @staticmethod
    def _write_csv(rows: List[List[str]], output_path: Path):
        """
        写入CSV文件

        Args:
            rows: 行数据列表
            output_path: 输出路径
        """
        with open(output_path, 'w', encoding='utf-8-sig', newline='') as f:
            writer = csv.writer(f)
            writer.writerows(rows)

    @staticmethod
    def export_extracted_data(
        data: Dict[str, Any],
        output_dir: str,
        formats: List[str] = None
    ) -> Dict[str, str]:
        """
        导出提取的数据到多种格式

        Args:
            data: 提取的文档数据
            output_dir: 输出目录
            formats: 导出格式列表,默认['json']。可选: 'json', 'csv'

        Returns:
            格式到文件路径的映射
        """
        if formats is None:
            formats = ['json']

        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        file_name = data.get("file_name", "unknown")
        base_name = Path(file_name).stem
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        results = {}

        for fmt in formats:
            if fmt == 'json':
                json_path = output_dir / f"{base_name}_{timestamp}.json"
                results['json'] = DataExporter.save_as_json(data, str(json_path))

            elif fmt == 'csv':
                if data.get("file_type") == "excel":
                    csv_path = output_dir / f"{base_name}_{timestamp}.csv"
                    results['csv'] = DataExporter.save_as_csv(data, str(csv_path))
                else:
                    print(f"⚠️  {file_name} 不是Excel文件,跳过CSV导出")

        return results
